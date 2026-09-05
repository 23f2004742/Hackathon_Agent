"""Security scanning, packaging, provenance, and the end-to-end pipeline."""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

import pytest

from hackathon_os import agents as roster
from hackathon_os.handoff import Status
from hackathon_os.llm import SimulatedBackend
from hackathon_os.orchestrator import Orchestrator
from hackathon_os.state import ProjectState
from hackathon_os.tools import REGISTRY, ExecutionContext, using


def ctx(agent: str, root: Path) -> ExecutionContext:
    spec = roster.get(agent)
    return ExecutionContext(
        root=root, agent=agent, write_paths=spec.write_paths,
        allowed_tools=frozenset(spec.tools), auto_approve=True,
    )


# -- secret scanning --------------------------------------------------------


def test_scanner_finds_a_planted_key(tmp_path):
    (tmp_path / "config.py").write_text(
        'ANTHROPIC_KEY = "sk-ant-api03-RealLookingKeyMaterial1234567890abcdef"\n',
        encoding="utf-8",
    )
    with using(ctx("security_reviewer", tmp_path)):
        out = REGISTRY["scan_secrets"].fn(path=".")
    assert "LIVE-LOOKING SECRETS" in out
    assert "Anthropic API key" in out


def test_scanner_separates_placeholders_from_leaks(tmp_path):
    (tmp_path / ".env.example").write_text(
        "ANTHROPIC_API_KEY=sk-ant-your-key-here-placeholder-value\n", encoding="utf-8"
    )
    with using(ctx("security_reviewer", tmp_path)):
        out = REGISTRY["scan_secrets"].fn(path=".")
    assert "PLACEHOLDERS" in out
    assert "LIVE-LOOKING SECRETS" not in out


def test_scanner_passes_a_clean_tree(tmp_path):
    (tmp_path / "app.py").write_text("import os\nkey = os.environ['API_KEY']\n", encoding="utf-8")
    with using(ctx("security_reviewer", tmp_path)):
        out = REGISTRY["scan_secrets"].fn(path=".")
    assert "PASS" in out


def test_code_scanner_flags_disabled_tls(tmp_path):
    (tmp_path / "client.py").write_text(
        "import requests\nrequests.get(url, verify=False)\n", encoding="utf-8"
    )
    with using(ctx("security_reviewer", tmp_path)):
        out = REGISTRY["scan_code_security"].fn(path=".")
    assert "TLS verification disabled" in out


# -- packaging --------------------------------------------------------------


def test_zip_never_packages_secrets(tmp_path):
    (tmp_path / "SUBMISSION").mkdir()
    (tmp_path / "README.md").write_text("# hi", encoding="utf-8")
    (tmp_path / ".env").write_text("ANTHROPIC_API_KEY=sk-ant-real-key-here", encoding="utf-8")
    with using(ctx("submission_manager", tmp_path)):
        out = REGISTRY["build_zip"].fn(
            output="SUBMISSION/submission.zip", include="**/*.md,**/*"
        )
    assert "EXCLUDED AS SECRETS" in out
    names = zipfile.ZipFile(tmp_path / "SUBMISSION/submission.zip").namelist()
    assert "README.md" in names
    assert not any(".env" in n for n in names)


# -- provenance -------------------------------------------------------------


def test_unsourced_claims_are_flagged(tmp_path):
    (tmp_path / "RESEARCH").mkdir()
    (tmp_path / "PRESENTATION").mkdir()
    (tmp_path / "PRESENTATION/pitch_strategy.md").write_text(
        "The market is worth $4.2B and we are 10x faster than anyone.\n", encoding="utf-8"
    )
    with using(ctx("pitch_strategist", tmp_path)):
        out = REGISTRY["verify_claims"].fn(document="PRESENTATION/pitch_strategy.md")
    assert "no matching ledger entry" in out


def test_recorded_claims_stop_being_flagged(tmp_path):
    (tmp_path / "RESEARCH").mkdir()
    (tmp_path / "PRESENTATION").mkdir()
    with using(ctx("market_researcher", tmp_path)):
        REGISTRY["record_source"].fn(
            claim="The triage software market is worth $4.2B.",
            url="https://example.org/report", title="Report",
        )
    (tmp_path / "PRESENTATION/pitch_strategy.md").write_text(
        "The market is worth $4.2B.\n", encoding="utf-8"
    )
    with using(ctx("pitch_strategist", tmp_path)):
        out = REGISTRY["verify_claims"].fn(document="PRESENTATION/pitch_strategy.md")
    assert "no unsourced factual claims" in out


def test_estimates_are_marked_as_unsourced(tmp_path):
    (tmp_path / "RESEARCH").mkdir()
    with using(ctx("market_researcher", tmp_path)):
        out = REGISTRY["record_source"].fn(claim="Our own estimate of TAM.", url="")
    assert "ESTIMATE" in out
    rows = json.loads((tmp_path / "RESEARCH/sources.json").read_text(encoding="utf-8"))
    assert rows[0]["kind"] == "estimate"


# -- end to end -------------------------------------------------------------


@pytest.fixture(scope="module")
def finished(tmp_path_factory) -> ProjectState:
    """One full autonomous run, on the simulated backend."""
    root = tmp_path_factory.mktemp("e2e") / "project"
    st = ProjectState.create(
        root, "E2E Test",
        problem=(
            "Build an AI-powered triage assistant for nurses, with a model that "
            "scores urgency, an interface on a tablet, and an auditable queue. "
            "Judged on impact, technical depth and whether a judge can run it."
        ),
        judging="Impact 30%, technical depth 25%, demo 20%, docs 10%.",
        submission="A zip named submission.zip, a README.md, and a .pptx deck.",
        constraints="Python only. Offline demo. 24 hours.",
    )
    orch = Orchestrator(st, SimulatedBackend(verbose=False), parallel=3, verbose=False)
    orch.plan()
    orch.run(max_waves=25)
    return st


def test_e2e_completes_every_task(finished):
    counts = finished.graph.counts()
    assert counts["failed"] == 0, [
        (t.id, t.result.notes) for t in finished.graph.tasks.values()
        if t.status is Status.FAILED
    ]
    assert counts["pending"] == 0
    assert counts["completed"] == counts["total"]


def test_e2e_reaches_the_done_phase(finished):
    assert finished.phase == "done"
    assert finished.graph.progress == 1.0


def test_e2e_produces_the_full_deliverable_set(finished):
    for rel in (
        "PRODUCT/requirements.md", "PRODUCT/product_plan.md", "PRODUCT/architecture.md",
        "RESEARCH/technical_research.md", "DESIGN/ux.md",
        "VALIDATION/test_report.md", "VALIDATION/security_review.md",
        "DEMO/demo_script.md", "README.md", "DOCUMENTATION/technical.md",
        "PRESENTATION/slides.md", "PRESENTATION/presentation.pptx",
        "FINAL/final_audit.md", "SUBMISSION/submission_manifest.json",
    ):
        p = finished.root / rel
        assert p.is_file(), f"missing deliverable: {rel}"
        assert p.stat().st_size > 100, f"stub deliverable: {rel}"


def test_e2e_pptx_is_a_real_openable_deck(finished):
    from pptx import Presentation
    prs = Presentation(str(finished.root / "PRESENTATION/presentation.pptx"))
    assert len(prs.slides) >= 6


def test_e2e_generated_python_actually_parses(finished):
    import ast
    for py in (finished.root / "src").rglob("*.py"):
        ast.parse(py.read_text(encoding="utf-8"))


def test_e2e_every_agent_reported_a_handoff(finished):
    assert len(finished.history) >= len(finished.graph.tasks)
    for r in finished.history:
        assert r.agent
        assert r.summary


def test_e2e_state_survives_reload(finished):
    reloaded = ProjectState.load(finished.root)
    assert reloaded.phase == "done"
    assert len(reloaded.graph.tasks) == len(finished.graph.tasks)
    assert len(reloaded.history) == len(finished.history)


def test_e2e_wrote_the_plan_and_decision_log(finished):
    assert (finished.root / "AGENT/plan.md").is_file()
    plan = (finished.root / "AGENT/plan.md").read_text(encoding="utf-8")
    assert "Deliberately not activated" in plan


def test_zip_keeps_env_example_but_drops_env(tmp_path):
    """.env.example documents required config; .env leaks it."""
    (tmp_path / "SUBMISSION").mkdir()
    (tmp_path / "README.md").write_text("# hi", encoding="utf-8")
    (tmp_path / ".env.example").write_text("API_KEY=your-key-here", encoding="utf-8")
    (tmp_path / ".env").write_text("API_KEY=sk-ant-real", encoding="utf-8")
    with using(ctx("submission_manager", tmp_path)):
        REGISTRY["build_zip"].fn(output="SUBMISSION/s.zip", include="**/*")
    names = zipfile.ZipFile(tmp_path / "SUBMISSION/s.zip").namelist()
    assert ".env.example" in names
    assert ".env" not in names
