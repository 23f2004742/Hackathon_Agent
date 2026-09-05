"""Document, presentation and packaging tools.

Real file generation: PPTX via python-pptx, PDF via reportlab, XLSX via
openpyxl, ZIP via stdlib. Each degrades to a clear error string naming the
missing library rather than silently producing nothing.
"""

from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from pathlib import Path

from .base import approve, fail, guard, tool, truncate


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def _parse_slides(md: str) -> list[tuple[str, list[str], str]]:
    """Split a slides markdown file into (title, bullets, notes) tuples.

    Format: `## Title` starts a slide, `-` lines are bullets, and a line
    beginning `Notes:` becomes speaker notes.
    """
    slides: list[tuple[str, list[str], str]] = []
    title, bullets, notes = None, [], []
    for raw in md.splitlines():
        line = raw.rstrip()
        if line.startswith("## "):
            if title is not None:
                slides.append((title, bullets, "\n".join(notes).strip()))
            title, bullets, notes = line[3:].strip(), [], []
        elif title is not None:
            s = line.strip()
            if s.lower().startswith("notes:"):
                notes.append(s[6:].strip())
            elif s.startswith(("- ", "* ")):
                bullets.append(s[2:].strip())
            elif s and not s.startswith("#"):
                notes.append(s)
    if title is not None:
        slides.append((title, bullets, "\n".join(notes).strip()))
    return slides


@tool("presentations", writes=True, approval=True)
def build_pptx(slides_markdown: str, output: str) -> str:
    """Render a slides markdown file into a real .pptx deck.

    Write the deck as markdown first, then call this. Format: `## Slide Title`
    per slide, `- bullet` for bullets, and `Notes: ...` for speaker notes.

    Args:
        slides_markdown: Path to the source markdown, relative to project root.
        output: Path for the .pptx to write, relative to project root.
    """
    try:
        ctx = guard("build_pptx")
        src = ctx.resolve(slides_markdown)
        if not src.is_file():
            return f"no such file: {slides_markdown}"
        dst = ctx.resolve_for_write(output)
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return "python-pptx is not installed; run: pip install python-pptx"

        slides = _parse_slides(src.read_text(encoding="utf-8", errors="replace"))
        if not slides:
            return f"{slides_markdown} contains no '## ' slide headings"
        if not approve(ctx, "build", f"{output} ({len(slides)} slides)"):
            return "User declined."
        if ctx.dry_run:
            return f"[dry-run] would build {output} with {len(slides)} slides"

        prs = Presentation()
        for i, (title, bullets, notes) in enumerate(slides):
            layout = prs.slide_layouts[0 if i == 0 and not bullets else 1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = title
            if bullets and len(s.placeholders) > 1:
                body = s.placeholders[1].text_frame
                body.text = bullets[0]
                for b in bullets[1:]:
                    para = body.add_paragraph()
                    para.text = b
                    para.level = 0
                for para in body.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(18)
            elif not bullets and len(s.placeholders) > 1:
                try:
                    s.placeholders[1].text = notes.splitlines()[0] if notes else ""
                except (KeyError, IndexError):
                    pass
            if notes:
                s.notes_slide.notes_text_frame.text = notes
        dst.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(dst))
        ctx.note(f"built {output}")
        return f"built {output}: {len(slides)} slides, {dst.stat().st_size} bytes"
    except Exception as e:  # noqa: BLE001
        return fail(e)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

@tool("documents", writes=True, approval=True)
def build_pdf(markdown_path: str, output: str, title: str = "") -> str:
    """Render a markdown document into a real .pdf.

    Handles headings, paragraphs, bullets and fenced code blocks. Tables are
    rendered as monospaced text rather than laid out.

    Args:
        markdown_path: Source markdown, relative to project root.
        output: Path for the .pdf to write, relative to project root.
        title: Document title for the PDF metadata. Defaults to the filename.
    """
    try:
        ctx = guard("build_pdf")
        src = ctx.resolve(markdown_path)
        if not src.is_file():
            return f"no such file: {markdown_path}"
        dst = ctx.resolve_for_write(output)
        try:
            from reportlab.lib.enums import TA_LEFT
            from reportlab.lib.pagesizes import LETTER
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import inch
            from reportlab.platypus import (
                ListFlowable, ListItem, PageBreak, Paragraph, Preformatted,
                SimpleDocTemplate, Spacer,
            )
        except ImportError:
            return "reportlab is not installed; run: pip install reportlab"

        if not approve(ctx, "build", output):
            return "User declined."
        if ctx.dry_run:
            return f"[dry-run] would build {output}"

        ss = getSampleStyleSheet()
        code_style = ParagraphStyle(
            "code", parent=ss["Code"], fontSize=8, leading=10, alignment=TA_LEFT
        )
        dst.parent.mkdir(parents=True, exist_ok=True)
        doc = SimpleDocTemplate(
            str(dst), pagesize=LETTER, title=title or dst.stem,
            leftMargin=0.9 * inch, rightMargin=0.9 * inch,
            topMargin=0.8 * inch, bottomMargin=0.8 * inch,
        )

        def esc(t: str) -> str:
            t = t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            t = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", t)
            t = re.sub(r"(?<!\*)\*([^*]+?)\*(?!\*)", r"<i>\1</i>", t)
            t = re.sub(r"`([^`]+?)`", r"<font face='Courier'>\1</font>", t)
            return t

        flow, bullets, in_code, code = [], [], False, []
        for raw in src.read_text(encoding="utf-8", errors="replace").splitlines():
            line = raw.rstrip()
            if line.strip().startswith("```"):
                if in_code:
                    flow.append(Preformatted("\n".join(code), code_style))
                    flow.append(Spacer(1, 8))
                    code = []
                in_code = not in_code
                continue
            if in_code:
                code.append(line)
                continue
            if bullets and not line.strip().startswith(("- ", "* ")):
                flow.append(ListFlowable(
                    [ListItem(Paragraph(esc(b), ss["BodyText"])) for b in bullets],
                    bulletType="bullet", leftIndent=18,
                ))
                flow.append(Spacer(1, 6))
                bullets = []
            s = line.strip()
            if not s:
                continue
            if s == "---":
                flow.append(PageBreak())
            elif s.startswith("#"):
                level = len(s) - len(s.lstrip("#"))
                flow.append(Spacer(1, 10))
                flow.append(Paragraph(esc(s.lstrip("# ").strip()), ss[f"Heading{min(level, 4)}"]))
            elif s.startswith(("- ", "* ")):
                bullets.append(s[2:])
            elif s.startswith("|"):
                flow.append(Preformatted(line, code_style))
            else:
                flow.append(Paragraph(esc(s), ss["BodyText"]))
        if bullets:
            flow.append(ListFlowable(
                [ListItem(Paragraph(esc(b), ss["BodyText"])) for b in bullets],
                bulletType="bullet", leftIndent=18,
            ))
        if in_code and code:
            flow.append(Preformatted("\n".join(code), code_style))
        if not flow:
            return f"{markdown_path} produced no renderable content"
        doc.build(flow)
        ctx.note(f"built {output}")
        return f"built {output}: {dst.stat().st_size} bytes"
    except Exception as e:  # noqa: BLE001
        return fail(e)


# ---------------------------------------------------------------------------
# Data
# ---------------------------------------------------------------------------

@tool("data", writes=True)
def write_csv(path: str, rows_json: str) -> str:
    """Write a CSV file from a JSON array of objects.

    Args:
        path: Output path relative to project root.
        rows_json: JSON array of flat objects; keys of the first become headers.
    """
    try:
        ctx = guard("write_csv")
        rows = json.loads(rows_json)
        if not isinstance(rows, list) or not rows:
            return "rows_json must be a non-empty JSON array of objects"
        p = ctx.resolve_for_write(path)
        buf = io.StringIO()
        w = csv.DictWriter(buf, fieldnames=list(rows[0].keys()))
        w.writeheader()
        w.writerows(rows)
        if not ctx.dry_run:
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(buf.getvalue(), encoding="utf-8", newline="")
        return f"wrote {path}: {len(rows)} rows"
    except Exception as e:  # noqa: BLE001
        return fail(e)


@tool("data")
def read_data(path: str, max_rows: int = 30) -> str:
    """Preview a CSV or JSON data file: shape, columns and the first rows.

    Args:
        path: Data file relative to project root.
        max_rows: How many rows to show. Default 30.
    """
    try:
        ctx = guard("read_data")
        p = ctx.resolve(path)
        if not p.is_file():
            return f"no such file: {path}"
        if p.suffix.lower() == ".json":
            return truncate(json.dumps(json.loads(p.read_text(encoding="utf-8")), indent=2))
        rows = list(csv.reader(io.StringIO(p.read_text(encoding="utf-8", errors="replace"))))
        if not rows:
            return f"{path} is empty"
        head = "\n".join(", ".join(r) for r in rows[: max_rows + 1])
        return truncate(f"[{path}: {len(rows) - 1} data rows x {len(rows[0])} cols]\n{head}")
    except Exception as e:  # noqa: BLE001
        return fail(e)


# ---------------------------------------------------------------------------
# Packaging
# ---------------------------------------------------------------------------

# Never let these into a submission archive, whatever the include globs say.
SECRET_PATTERNS = (
    ".env", ".env.*", "*.pem", "*.key", "id_rsa*", "*.p12", "*credentials*",
    "*secret*", "*.sqlite", ".git/*",
)

# Template files that look like secrets but carry no secret and belong in the
# submission -- they are how a judge learns what configuration to supply.
SECRET_EXCEPTIONS = (".env.example", ".env.sample", ".env.template")


@tool("packaging", writes=True, approval=True)
def build_zip(output: str, include: str = "**/*", exclude: str = "") -> str:
    """Package project files into a .zip for submission.

    Secrets are excluded unconditionally -- .env files, keys and credentials are
    dropped even if your include glob matches them, and the result reports what
    was skipped. Verify that report before submitting.

    Args:
        output: Path for the .zip, relative to project root.
        include: Comma-separated globs to include. Default "**/*".
        exclude: Comma-separated globs to exclude, in addition to secrets.
    """
    try:
        ctx = guard("build_zip")
        dst = ctx.resolve_for_write(output)
        incs = [g.strip() for g in include.split(",") if g.strip()]
        excs = [g.strip() for g in exclude.split(",") if g.strip()]

        chosen: list[Path] = []
        for g in incs:
            chosen.extend(p for p in ctx.root.glob(g) if p.is_file())
        chosen = sorted(set(chosen))

        keep, skipped_secret, skipped_user = [], [], []
        for p in chosen:
            rel = p.relative_to(ctx.root).as_posix()
            if p.resolve() == dst.resolve():
                continue
            name = Path(rel).name
            if name not in SECRET_EXCEPTIONS and any(
                Path(rel).match(s) or name == s or rel.startswith(".git/")
                for s in SECRET_PATTERNS
            ):
                skipped_secret.append(rel)
                continue
            if any(part in {"__pycache__", ".git", "node_modules", ".venv"}
                   for part in Path(rel).parts):
                continue
            if any(Path(rel).match(x) for x in excs):
                skipped_user.append(rel)
                continue
            keep.append((p, rel))

        if not keep:
            return "nothing to package: include glob matched no eligible files"
        if not approve(ctx, "package", f"{output} ({len(keep)} files)"):
            return "User declined."
        if ctx.dry_run:
            return f"[dry-run] would package {len(keep)} files into {output}"

        dst.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as z:
            for p, rel in keep:
                z.write(p, rel)
        size = dst.stat().st_size
        msg = f"packaged {output}: {len(keep)} files, {size / 1024:.0f} KB"
        if skipped_secret:
            msg += f"\n  EXCLUDED AS SECRETS ({len(skipped_secret)}): " + ", ".join(skipped_secret[:10])
        if skipped_user:
            msg += f"\n  excluded by your globs ({len(skipped_user)})"
        ctx.note(f"packaged {output} ({len(keep)} files)")
        return msg
    except Exception as e:  # noqa: BLE001
        return fail(e)
